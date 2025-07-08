"""
テキスト抽出モジュール
様々なドキュメント形式からテキストを抽出する
"""
import logging
from pathlib import Path
from typing import Optional

import PyPDF2
import fitz  # pymupdf
from docx import Document
import openpyxl
from pptx import Presentation

logger = logging.getLogger(__name__)


class TextExtractor:
    """テキスト抽出クラス"""
    
    @staticmethod
    def extract_text(file_path: str) -> str:
        """
        ファイルからテキストを抽出する
        
        Args:
            file_path: 抽出対象のファイルパス
            
        Returns:
            抽出されたテキスト
            
        Raises:
            Exception: テキスト抽出に失敗した場合
        """
        try:
            logger.info(f"テキスト抽出を開始: {file_path}")
            
            # ファイル拡張子に基づく処理
            file_extension = Path(file_path).suffix.lower()
            
            if file_extension in ['.pdf']:
                return TextExtractor._extract_from_pdf(file_path)
            elif file_extension in ['.docx', '.doc']:
                return TextExtractor._extract_from_word(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                return TextExtractor._extract_from_excel(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                return TextExtractor._extract_from_powerpoint(file_path)
            else:
                raise Exception(f"サポートされていないファイル形式です: {file_extension}")
                
        except Exception as e:
            logger.error(f"テキスト抽出に失敗しました: {file_path}, エラー: {e}")
            raise Exception(f"テキスト抽出に失敗しました: {str(e)}")
    
    @staticmethod
    def _extract_from_pdf(file_path: str) -> str:
        """PDFからテキストを抽出"""
        try:
            # まずPyMuPDFを試す（高性能でOCR機能もある）
            try:
                doc = fitz.open(file_path)
                text = ""
                for page in doc:
                    page_text = page.get_text()
                    if page_text:
                        text += page_text + "\n"
                doc.close()
                if text.strip():
                    logger.info(f"PyMuPDFでPDFテキスト抽出成功: {len(text)}文字")
                    return text
            except Exception as e:
                logger.warning(f"PyMuPDFでの抽出に失敗、PyPDF2を試行: {e}")
            
            # PyMuPDFが失敗した場合はPyPDF2を使用
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in pdf_reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    if text.strip():
                        logger.info(f"PyPDF2でPDFテキスト抽出成功: {len(text)}文字")
                        return text
            except Exception as e:
                logger.error(f"PyPDF2での抽出も失敗: {e}")
            
            raise Exception("すべてのPDF抽出ライブラリで処理に失敗しました")
            
        except Exception as e:
            logger.error(f"PDF抽出に失敗: {e}")
            raise Exception(f"PDFファイルのテキスト抽出に失敗しました: {str(e)}")
            
            raise Exception("すべてのPDF抽出方法が失敗しました")
            
        except Exception as e:
            logger.error(f"PDF抽出に失敗: {e}")
            raise
    
    @staticmethod
    def _extract_from_word(file_path: str) -> str:
        """Wordドキュメントからテキストを抽出"""
        try:
            # .docxファイルの場合はpython-docxを使用
            if file_path.lower().endswith('.docx'):
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            else:
                # .docファイルは現在サポートしていない
                raise Exception(".docファイルはサポートされていません。.docxファイルを使用してください。")
            
        except Exception as e:
            logger.error(f"Word抽出に失敗: {e}")
            raise
    
    @staticmethod
    def _extract_from_excel(file_path: str) -> str:
        """Excelファイルからテキストを抽出"""
        try:
            # .xlsxファイルの場合はopenpyxlを使用
            if file_path.lower().endswith('.xlsx'):
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                text = ""
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text += f"シート: {sheet_name}\n"
                    for row in sheet.iter_rows(values_only=True):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                row_text.append(str(cell))
                        if row_text:
                            text += "\t".join(row_text) + "\n"
                    text += "\n"
                return text
            else:
                # .xlsファイルは現在サポートしていない
                raise Exception(".xlsファイルはサポートされていません。.xlsxファイルを使用してください。")
            
        except Exception as e:
            logger.error(f"Excel抽出に失敗: {e}")
            raise
    
    @staticmethod
    def _extract_from_powerpoint(file_path: str) -> str:
        """PowerPointファイルからテキストを抽出"""
        try:
            # .pptxファイルの場合はpython-pptxを使用
            if file_path.lower().endswith('.pptx'):
                presentation = Presentation(file_path)
                text = ""
                for slide_num, slide in enumerate(presentation.slides, 1):
                    text += f"スライド {slide_num}:\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + "\n"
                    text += "\n"
                return text
            else:
                # .pptファイルは現在サポートしていない
                raise Exception(".pptファイルはサポートされていません。.pptxファイルを使用してください。")
            
        except Exception as e:
            logger.error(f"PowerPoint抽出に失敗: {e}")
            raise
    
    @staticmethod
    def clean_extracted_text(text: str) -> str:
        """
        抽出されたテキストをクリーンアップ
        
        Args:
            text: 抽出されたテキスト
            
        Returns:
            クリーンアップされたテキスト
        """
        if not text:
            return ""
        
        # 不要な空白や改行を整理
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            # 空白のみの行をスキップ
            line = line.strip()
            if line:
                cleaned_lines.append(line)
        
        # 連続する空行を削除
        result = '\n'.join(cleaned_lines)
        
        # 異常に長い連続するスペースを削除
        import re
        result = re.sub(r' {3,}', ' ', result)
        
        return result
    
    @staticmethod
    def get_text_stats(text: str) -> dict:
        """
        テキストの統計情報を取得
        
        Args:
            text: 対象テキスト
            
        Returns:
            統計情報の辞書
        """
        if not text:
            return {
                'character_count': 0,
                'line_count': 0,
                'word_count': 0,
                'paragraph_count': 0
            }
        
        lines = text.split('\n')
        non_empty_lines = [line for line in lines if line.strip()]
        
        # 段落数（空行で区切られたブロック）
        paragraphs = text.split('\n\n')
        non_empty_paragraphs = [p for p in paragraphs if p.strip()]
        
        # 単語数（日本語の場合は文字数ベース）
        import re
        words = re.findall(r'[ぁ-んァ-ヶ一-龯]+|[a-zA-Z]+', text)
        
        return {
            'character_count': len(text),
            'line_count': len(non_empty_lines),
            'word_count': len(words),
            'paragraph_count': len(non_empty_paragraphs)
        }
